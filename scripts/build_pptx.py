#!/usr/bin/env python3
"""
Build a real .pptx deck from factory slides.json.

Features:
  - Multiple themes
  - Schema validation
  - Error recovery (--recover)
  - Clear exit codes and logging

Usage:
  python scripts/build_pptx.py --input path/slides.json --output path/slides.pptx
  python scripts/build_pptx.py -i slides.json -o slides.pptx --theme academic-light --recover
"""

from __future__ import annotations

import argparse
import json
import logging
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError as exc:  # pragma: no cover
    print(
        "ERROR: python-pptx is required. Install with: pip install python-pptx",
        file=sys.stderr,
    )
    raise SystemExit(2) from exc


LOG = logging.getLogger("build_pptx")

SUPPORTED_LAYOUTS = frozenset(
    {"title", "section", "bullets", "two_column", "objectives", "closing"}
)

THEMES: dict[str, dict[str, Any]] = {
    "academic-light": {
        "bg": RGBColor(0xF7, 0xF5, 0xF0),
        "title": RGBColor(0x1B, 0x2A, 0x4A),
        "body": RGBColor(0x2C, 0x2C, 0x2C),
        "accent": RGBColor(0x0B, 0x6E, 0x4F),
        "muted": RGBColor(0x5C, 0x6B, 0x7A),
        "section_bg": RGBColor(0x1B, 0x2A, 0x4A),
        "section_fg": RGBColor(0xF7, 0xF5, 0xF0),
    },
    "academic-dark": {
        "bg": RGBColor(0x14, 0x18, 0x1F),
        "title": RGBColor(0xE8, 0xEE, 0xF7),
        "body": RGBColor(0xC5, 0xCD, 0xD8),
        "accent": RGBColor(0x3D, 0xB8, 0x8C),
        "muted": RGBColor(0x8A, 0x95, 0xA5),
        "section_bg": RGBColor(0x0E, 0x12, 0x18),
        "section_fg": RGBColor(0xE8, 0xEE, 0xF7),
    },
    "minimal-mono": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x11, 0x11, 0x11),
        "body": RGBColor(0x22, 0x22, 0x22),
        "accent": RGBColor(0x44, 0x44, 0x44),
        "muted": RGBColor(0x66, 0x66, 0x66),
        "section_bg": RGBColor(0x11, 0x11, 0x11),
        "section_fg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "campus-blue": {
        "bg": RGBColor(0xF3, 0xF7, 0xFB),
        "title": RGBColor(0x0A, 0x2E, 0x5C),
        "body": RGBColor(0x1F, 0x2A, 0x37),
        "accent": RGBColor(0x1F, 0x6F, 0xB2),
        "muted": RGBColor(0x4F, 0x64, 0x7A),
        "section_bg": RGBColor(0x0A, 0x2E, 0x5C),
        "section_fg": RGBColor(0xF3, 0xF7, 0xFB),
    },
}

FALLBACK_THEME = "minimal-mono"


class BuildError(Exception):
    """Unrecoverable build failure."""


def setup_logging(verbose: bool) -> None:
    level = logging.DEBUG if verbose else logging.INFO
    logging.basicConfig(
        level=level,
        format="[build_pptx] %(levelname)s: %(message)s",
    )


def load_json(path: Path) -> dict[str, Any]:
    try:
        raw = path.read_text(encoding="utf-8")
    except OSError as exc:
        raise BuildError(f"Cannot read input: {path}") from exc
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise BuildError(f"Invalid JSON in {path}: {exc}") from exc
    if not isinstance(data, dict):
        raise BuildError("Root JSON value must be an object")
    return data


def _as_str(value: Any, default: str = "") -> str:
    if value is None:
        return default
    return str(value).strip()


def _as_bullets(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        text = value.strip()
        return [text] if text else []
    if isinstance(value, list):
        out: list[str] = []
        for item in value:
            if item is None:
                continue
            if isinstance(item, dict):
                text = _as_str(item.get("text") or item.get("content"))
            else:
                text = _as_str(item)
            if text:
                out.append(text)
        return out
    return [_as_str(value)]


def normalize_deck(data: dict[str, Any], recover: bool, strict: bool) -> dict[str, Any]:
    meta_in = data.get("meta") or {}
    if not isinstance(meta_in, dict):
        if recover:
            LOG.warning("meta is not an object; substituting defaults")
            meta_in = {}
        else:
            raise BuildError("meta must be an object")

    theme = _as_str(meta_in.get("theme"), "academic-light") or "academic-light"
    if theme not in THEMES:
        msg = f"Unknown theme '{theme}'"
        if recover or theme:
            LOG.warning("%s — falling back to %s", msg, FALLBACK_THEME)
            theme = FALLBACK_THEME
        else:
            raise BuildError(msg)

    meta = {
        "title": _as_str(meta_in.get("title"), "Untitled Deck"),
        "subtitle": _as_str(meta_in.get("subtitle")),
        "author": _as_str(meta_in.get("author"), "Content Factory"),
        "subject_slug": _as_str(meta_in.get("subject_slug")),
        "version": _as_str(meta_in.get("version")),
        "theme": theme,
        "language": _as_str(meta_in.get("language"), "en") or "en",
    }

    slides_in = data.get("slides")
    if not isinstance(slides_in, list) or not slides_in:
        raise BuildError("slides must be a non-empty array")

    slides: list[dict[str, Any]] = []
    for index, raw in enumerate(slides_in, start=1):
        try:
            slide = normalize_slide(raw, index, recover=recover, strict=strict)
        except BuildError as exc:
            if recover:
                LOG.warning("Skipping slide %s: %s", index, exc)
                continue
            raise
        if slide is not None:
            slides.append(slide)

    if not slides:
        raise BuildError("No valid slides remain after normalization")

    # Ensure a title slide exists when recovering a broken deck
    if recover and slides[0]["layout"] != "title":
        LOG.warning("Prepending synthetic title slide")
        slides.insert(
            0,
            {
                "id": "s00",
                "layout": "title",
                "title": meta["title"],
                "subtitle": meta["subtitle"],
                "bullets": [],
                "notes": "",
                "left": {"heading": "", "bullets": []},
                "right": {"heading": "", "bullets": []},
            },
        )

    return {"meta": meta, "slides": slides}


def normalize_slide(
    raw: Any, index: int, recover: bool, strict: bool
) -> dict[str, Any] | None:
    if not isinstance(raw, dict):
        raise BuildError(f"slide {index} is not an object")

    if strict:
        allowed = {
            "id",
            "layout",
            "title",
            "subtitle",
            "bullets",
            "notes",
            "left",
            "right",
        }
        unknown = set(raw) - allowed
        if unknown:
            raise BuildError(f"slide {index} has unknown keys: {sorted(unknown)}")

    layout = _as_str(raw.get("layout"), "bullets") or "bullets"
    if layout not in SUPPORTED_LAYOUTS:
        msg = f"unsupported layout '{layout}'"
        if recover:
            LOG.warning("slide %s: %s — coercing to bullets", index, msg)
            layout = "bullets"
        else:
            raise BuildError(msg)

    title = _as_str(raw.get("title"))
    if not title:
        if recover:
            title = f"Slide {index}"
            LOG.warning("slide %s: empty title — using '%s'", index, title)
        else:
            raise BuildError(f"slide {index} missing title")

    slide_id = _as_str(raw.get("id"), f"s{index:02d}") or f"s{index:02d}"

    left = raw.get("left") if isinstance(raw.get("left"), dict) else {}
    right = raw.get("right") if isinstance(raw.get("right"), dict) else {}

    return {
        "id": slide_id,
        "layout": layout,
        "title": title,
        "subtitle": _as_str(raw.get("subtitle")),
        "bullets": _as_bullets(raw.get("bullets")),
        "notes": _as_str(raw.get("notes")),
        "left": {
            "heading": _as_str(left.get("heading")),
            "bullets": _as_bullets(left.get("bullets")),
        },
        "right": {
            "heading": _as_str(right.get("heading")),
            "bullets": _as_bullets(right.get("bullets")),
        },
    }


def set_run_font(run, *, size_pt: float, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"


def fill_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def write_title(tf, text: str, theme: dict[str, Any], size: float = 32, bold: bool = True) -> None:
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_font(run, size_pt=size, color=theme["title"], bold=bold)


def write_paragraphs(
    tf,
    lines: list[str],
    theme: dict[str, Any],
    *,
    size: float = 18,
    bullet: bool = True,
    color: RGBColor | None = None,
) -> None:
    tf.clear()
    fg = color or theme["body"]
    if not lines:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = ""
        set_run_font(run, size_pt=size, color=fg)
        return

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        if bullet:
            p.text = ""
            run = p.add_run()
            run.text = line
        else:
            run = p.add_run()
            run.text = line
        set_run_font(run, size_pt=size, color=fg)
        p.space_after = Pt(8)


def set_speaker_notes(slide, text: str) -> None:
    if not text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def render_title_slide(prs, slide_data: dict[str, Any], theme: dict[str, Any]) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    fill_background(slide, theme["bg"])

    accent = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        Inches(0.25),
        prs.slide_height,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme["accent"]
    accent.line.fill.background()

    box = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(12), Inches(1.5))
    write_title(box.text_frame, slide_data["title"], theme, size=36)

    if slide_data.get("subtitle"):
        sub = add_textbox(slide, Inches(0.8), Inches(3.7), Inches(12), Inches(1))
        write_paragraphs(
            sub.text_frame,
            [slide_data["subtitle"]],
            theme,
            size=18,
            bullet=False,
            color=theme["muted"],
        )

    set_speaker_notes(slide, slide_data.get("notes", ""))


def render_section_slide(prs, slide_data: dict[str, Any], theme: dict[str, Any]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    fill_background(slide, theme["section_bg"])

    box = add_textbox(slide, Inches(0.8), Inches(2.6), Inches(12), Inches(1.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data["title"]
    set_run_font(run, size_pt=34, color=theme["section_fg"], bold=True)

    set_speaker_notes(slide, slide_data.get("notes", ""))


def render_bullets_slide(
    prs,
    slide_data: dict[str, Any],
    theme: dict[str, Any],
    *,
    accent_label: str | None = None,
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    fill_background(slide, theme["bg"])

    if accent_label:
        tag = add_textbox(slide, Inches(0.7), Inches(0.35), Inches(3), Inches(0.4))
        write_paragraphs(
            tag.text_frame,
            [accent_label],
            theme,
            size=12,
            bullet=False,
            color=theme["accent"],
        )

    title_box = add_textbox(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(1))
    write_title(title_box.text_frame, slide_data["title"], theme, size=28)

    body = add_textbox(slide, Inches(0.9), Inches(1.8), Inches(11.5), Inches(5))
    bullets = slide_data.get("bullets") or []
    write_paragraphs(body.text_frame, bullets, theme, size=18, bullet=True)

    # Prefixed bullets for readability in PowerPoint
    for i, p in enumerate(body.text_frame.paragraphs):
        if i < len(bullets) and bullets[i]:
            # Rebuild with bullet character for portability
            text = bullets[i]
            p.clear()
            run = p.add_run()
            run.text = f"• {text}"
            set_run_font(run, size_pt=18, color=theme["body"])

    set_speaker_notes(slide, slide_data.get("notes", ""))


def render_two_column_slide(prs, slide_data: dict[str, Any], theme: dict[str, Any]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    fill_background(slide, theme["bg"])

    title_box = add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.9))
    write_title(title_box.text_frame, slide_data["title"], theme, size=28)

    left = slide_data.get("left") or {}
    right = slide_data.get("right") or {}

    left_box = add_textbox(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5))
    right_box = add_textbox(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5))

    def fill_column(tf, heading: str, bullets: list[str]) -> None:
        lines: list[str] = []
        if heading:
            lines.append(heading)
        lines.extend(f"• {b}" for b in bullets)
        write_paragraphs(tf, lines, theme, size=16, bullet=False)
        if heading and tf.paragraphs:
            # Bold the heading run
            p0 = tf.paragraphs[0]
            if p0.runs:
                set_run_font(p0.runs[0], size_pt=18, color=theme["accent"], bold=True)

    fill_column(left_box.text_frame, left.get("heading", ""), left.get("bullets") or [])
    fill_column(right_box.text_frame, right.get("heading", ""), right.get("bullets") or [])
    set_speaker_notes(slide, slide_data.get("notes", ""))


def render_slide(prs, slide_data: dict[str, Any], theme: dict[str, Any]) -> None:
    layout = slide_data["layout"]
    if layout == "title":
        render_title_slide(prs, slide_data, theme)
    elif layout == "section":
        render_section_slide(prs, slide_data, theme)
    elif layout == "two_column":
        render_two_column_slide(prs, slide_data, theme)
    elif layout == "objectives":
        render_bullets_slide(prs, slide_data, theme, accent_label="OBJECTIVES")
    elif layout == "closing":
        render_bullets_slide(prs, slide_data, theme, accent_label="CLOSE")
    else:
        render_bullets_slide(prs, slide_data, theme)


def build_presentation(deck: dict[str, Any]) -> Presentation:
    theme_name = deck["meta"]["theme"]
    theme = THEMES[theme_name]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in deck["slides"]:
        try:
            render_slide(prs, slide_data, theme)
        except Exception as exc:  # noqa: BLE001 — per-slide isolation
            raise BuildError(
                f"Failed rendering slide {slide_data.get('id')}: {exc}"
            ) from exc
    return prs


def build_minimal_fallback(meta: dict[str, Any], raw_slides: list[Any]) -> dict[str, Any]:
    """Last-resort deck: title, objectives-ish bullets, topic titles, closing."""
    LOG.warning("Building minimal fallback deck")
    titles: list[str] = []
    for item in raw_slides:
        if isinstance(item, dict):
            t = _as_str(item.get("title"))
            if t:
                titles.append(t)
    titles = titles[:12] or ["Overview", "Key ideas", "Practice", "Summary"]

    slides = [
        {
            "id": "s01",
            "layout": "title",
            "title": meta.get("title") or "Untitled Deck",
            "subtitle": meta.get("subtitle") or "",
            "bullets": [],
            "notes": "Minimal recovery deck.",
            "left": {"heading": "", "bullets": []},
            "right": {"heading": "", "bullets": []},
        },
        {
            "id": "s02",
            "layout": "objectives",
            "title": "Learning objectives",
            "subtitle": "",
            "bullets": ["Review the outline objectives with students"],
            "notes": "",
            "left": {"heading": "", "bullets": []},
            "right": {"heading": "", "bullets": []},
        },
    ]
    for i, title in enumerate(titles, start=3):
        slides.append(
            {
                "id": f"s{i:02d}",
                "layout": "bullets",
                "title": title,
                "subtitle": "",
                "bullets": ["Key point 1", "Key point 2", "Discussion prompt"],
                "notes": "",
                "left": {"heading": "", "bullets": []},
                "right": {"heading": "", "bullets": []},
            }
        )
    slides.append(
        {
            "id": f"s{len(slides)+1:02d}",
            "layout": "closing",
            "title": "Questions & next steps",
            "subtitle": "",
            "bullets": ["Ask clarifying questions", "Complete practice materials"],
            "notes": "",
            "left": {"heading": "", "bullets": []},
            "right": {"heading": "", "bullets": []},
        }
    )
    meta = {**meta, "theme": meta.get("theme") if meta.get("theme") in THEMES else FALLBACK_THEME}
    return {"meta": meta, "slides": slides}


def save_presentation(prs: Presentation, output: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    tmp = output.with_suffix(output.suffix + ".tmp")
    try:
        prs.save(str(tmp))
        tmp.replace(output)
    except OSError as exc:
        if tmp.exists():
            try:
                tmp.unlink()
            except OSError:
                pass
        raise BuildError(f"Failed to write {output}: {exc}") from exc


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Build PPTX from content-factory slides.json"
    )
    parser.add_argument("--input", "-i", required=True, type=Path, help="Path to slides.json")
    parser.add_argument("--output", "-o", required=True, type=Path, help="Path to slides.pptx")
    parser.add_argument(
        "--theme",
        "-t",
        default=None,
        choices=sorted(THEMES.keys()),
        help="Override theme from JSON",
    )
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Fail on unknown keys / empty titles",
    )
    parser.add_argument(
        "--recover",
        action="store_true",
        help="Enable coercion, skip bad slides, theme fallback",
    )
    parser.add_argument(
        "--minimal-on-fail",
        action="store_true",
        help="If normal build fails, emit a minimal recovery deck",
    )
    parser.add_argument("--verbose", "-v", action="store_true")
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv)
    setup_logging(args.verbose)

    try:
        raw = load_json(args.input)
    except BuildError as exc:
        LOG.error("%s", exc)
        return 1

    recover = bool(args.recover)
    strict = bool(args.strict)

    try:
        deck = normalize_deck(raw, recover=recover, strict=strict)
        if args.theme:
            if args.theme not in THEMES:
                LOG.warning("CLI theme invalid; using %s", FALLBACK_THEME)
                deck["meta"]["theme"] = FALLBACK_THEME
            else:
                deck["meta"]["theme"] = args.theme
        prs = build_presentation(deck)
        save_presentation(prs, args.output)
    except BuildError as exc:
        LOG.error("%s", exc)
        if args.minimal_on_fail or recover:
            try:
                meta = raw.get("meta") if isinstance(raw.get("meta"), dict) else {}
                meta = {
                    "title": _as_str(meta.get("title"), "Untitled Deck"),
                    "subtitle": _as_str(meta.get("subtitle")),
                    "author": _as_str(meta.get("author"), "Content Factory"),
                    "theme": FALLBACK_THEME,
                }
                if args.theme in THEMES:
                    meta["theme"] = args.theme
                deck = build_minimal_fallback(meta, raw.get("slides") or [])
                prs = build_presentation(deck)
                save_presentation(prs, args.output)
                LOG.warning("Wrote minimal recovery PPTX to %s", args.output)
                return 0
            except Exception as nested:  # noqa: BLE001
                LOG.error("Recovery failed: %s", nested)
                return 1
        return 1
    except Exception as exc:  # noqa: BLE001
        LOG.exception("Unexpected failure: %s", exc)
        return 1

    LOG.info(
        "Wrote %s (%s slides, theme=%s)",
        args.output,
        len(deck["slides"]),
        deck["meta"]["theme"],
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
